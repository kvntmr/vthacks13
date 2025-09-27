"""
PowerPoint Parser Service
Extracts text from PowerPoint files and uses OCR for non-text elements
"""

import os
import io
import tempfile
from typing import Dict, List, Any, Optional, Tuple
from pathlib import Path
import asyncio
import aiofiles

# PowerPoint processing
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Image processing and OCR
from PIL import Image
import pytesseract
import cv2
import numpy as np

class PowerPointParser:
    """Parser for PowerPoint files with text extraction and OCR fallback"""
    
    def __init__(self, tesseract_path: str = ""):
        """
        Initialize the PowerPoint parser
        
        Args:
            tesseract_path: Path to tesseract executable (if not in PATH, leave empty to use system PATH)
        """
        if tesseract_path:
            pytesseract.pytesseract.tesseract_cmd = tesseract_path
        
        # OCR configuration - simplified to avoid quote issues
        self.ocr_config = '--oem 3 --psm 6'
    
    async def parse_powerpoint(self, file_path: str) -> Dict[str, Any]:
        """
        Parse a PowerPoint file and extract all text content
        
        Args:
            file_path: Path to the PowerPoint file
            
        Returns:
            Dictionary containing extracted text and metadata
        """
        try:
            # Load the presentation
            presentation = Presentation(file_path)
            
            # Initialize result structure
            result = {
                "file_path": file_path,
                "file_name": os.path.basename(file_path),
                "total_slides": len(presentation.slides),
                "slides": [],
                "extracted_text": "",
                "ocr_used": False,
                "processing_summary": {
                    "text_boxes": 0,
                    "images_processed": 0,
                    "tables_processed": 0,
                    "shapes_processed": 0
                }
            }
            
            # Process each slide
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_data = await self._process_slide(slide, slide_num)
                result["slides"].append(slide_data)
                
                # Update processing summary
                result["processing_summary"]["text_boxes"] += slide_data["text_box_count"]
                result["processing_summary"]["images_processed"] += slide_data["images_processed"]
                result["processing_summary"]["tables_processed"] += slide_data["tables_processed"]
                result["processing_summary"]["shapes_processed"] += slide_data["shapes_processed"]
                
                # Collect all text
                result["extracted_text"] += slide_data["slide_text"] + "\n"
            
            # Clean up extracted text
            result["extracted_text"] = result["extracted_text"].strip()
            
            return result
            
        except Exception as e:
            return {
                "error": f"Failed to parse PowerPoint file: {str(e)}",
                "file_path": file_path,
                "file_name": os.path.basename(file_path)
            }
    
    async def _process_slide(self, slide, slide_num: int) -> Dict[str, Any]:
        """
        Process a single slide and extract all text content
        
        Args:
            slide: PowerPoint slide object
            slide_num: Slide number
            
        Returns:
            Dictionary containing slide data and extracted text
        """
        slide_data = {
            "slide_number": slide_num,
            "slide_text": "",
            "text_boxes": [],
            "images": [],
            "tables": [],
            "shapes": [],
            "text_box_count": 0,
            "images_processed": 0,
            "tables_processed": 0,
            "shapes_processed": 0,
            "ocr_used": False
        }
        
        # Process all shapes in the slide
        for shape in slide.shapes:
            shape_data = await self._process_shape(shape)
            
            if shape_data["type"] == "text":
                slide_data["text_boxes"].append(shape_data)
                slide_data["text_box_count"] += 1
                slide_data["slide_text"] += shape_data["text"] + "\n"
                
            elif shape_data["type"] == "image":
                slide_data["images"].append(shape_data)
                slide_data["images_processed"] += 1
                if shape_data["ocr_text"]:
                    slide_data["slide_text"] += shape_data["ocr_text"] + "\n"
                    slide_data["ocr_used"] = True
                    
            elif shape_data["type"] == "table":
                slide_data["tables"].append(shape_data)
                slide_data["tables_processed"] += 1
                slide_data["slide_text"] += shape_data["text"] + "\n"
                
            elif shape_data["type"] == "shape":
                slide_data["shapes"].append(shape_data)
                slide_data["shapes_processed"] += 1
                if shape_data["text"]:
                    slide_data["slide_text"] += shape_data["text"] + "\n"
        
        # Clean up slide text
        slide_data["slide_text"] = slide_data["slide_text"].strip()
        
        return slide_data
    
    async def _process_shape(self, shape) -> Dict[str, Any]:
        """
        Process a single shape and extract text content
        
        Args:
            shape: PowerPoint shape object
            
        Returns:
            Dictionary containing shape data and extracted text
        """
        shape_data = {
            "type": "unknown",
            "text": "",
            "ocr_text": "",
            "position": {
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height
            }
        }
        
        try:
            # Handle different shape types
            if shape.has_text_frame:
                # Text box or shape with text
                shape_data["type"] = "text"
                shape_data["text"] = self._extract_text_from_frame(shape.text_frame)
                
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Image - use OCR
                shape_data["type"] = "image"
                shape_data["ocr_text"] = await self._extract_text_from_image(shape)
                
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                # Table
                shape_data["type"] = "table"
                shape_data["text"] = self._extract_text_from_table(shape.table)
                
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # Group of shapes - process recursively
                shape_data["type"] = "group"
                group_text = ""
                for sub_shape in shape.shapes:
                    sub_data = await self._process_shape(sub_shape)
                    if sub_data["text"]:
                        group_text += sub_data["text"] + "\n"
                    if sub_data["ocr_text"]:
                        group_text += sub_data["ocr_text"] + "\n"
                shape_data["text"] = group_text.strip()
                
            else:
                # Other shapes
                shape_data["type"] = "shape"
                if hasattr(shape, 'text') and shape.text:
                    shape_data["text"] = shape.text
                    
        except Exception as e:
            shape_data["error"] = f"Error processing shape: {str(e)}"
        
        return shape_data
    
    def _extract_text_from_frame(self, text_frame) -> str:
        """Extract text from a text frame"""
        text = ""
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                text += run.text
            text += "\n"
        return text.strip()
    
    def _extract_text_from_table(self, table) -> str:
        """Extract text from a table"""
        text = ""
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                cell_text = ""
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        cell_text += run.text
                row_text.append(cell_text.strip())
            text += " | ".join(row_text) + "\n"
        return text.strip()
    
    async def _extract_text_from_image(self, shape) -> str:
        """
        Extract text from an image using OCR
        
        Args:
            shape: PowerPoint shape containing an image
            
        Returns:
            Extracted text from the image
        """
        try:
            # Get image data from shape
            image_data = shape.image.blob
            
            # Convert to PIL Image
            image = Image.open(io.BytesIO(image_data))
            
            # Convert to OpenCV format for preprocessing
            cv_image = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2BGR)
            
            # Preprocess image for better OCR
            processed_image = self._preprocess_image_for_ocr(cv_image)
            
            # Try to perform OCR with error handling
            try:
                text = pytesseract.image_to_string(processed_image, config=self.ocr_config)
                return text.strip()
            except Exception as ocr_error:
                # If OCR fails due to missing language data, provide helpful message
                if "tessdata" in str(ocr_error).lower() or "language" in str(ocr_error).lower():
                    return "[OCR Error: Tesseract language data not found. Please install tesseract-data-eng package or download eng.traineddata]"
                else:
                    return f"[OCR Error: {str(ocr_error)}]"
            
        except Exception as e:
            return f"[Image Processing Error: {str(e)}]"
    
    def _preprocess_image_for_ocr(self, image: np.ndarray) -> np.ndarray:
        """
        Preprocess image for better OCR results
        
        Args:
            image: OpenCV image array
            
        Returns:
            Preprocessed image array
        """
        try:
            # Convert to grayscale
            gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
            
            # Apply Gaussian blur to reduce noise
            blurred = cv2.GaussianBlur(gray, (5, 5), 0)
            
            # Apply threshold to get binary image
            _, thresh = cv2.threshold(blurred, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
            
            # Morphological operations to clean up
            kernel = np.ones((2, 2), np.uint8)
            cleaned = cv2.morphologyEx(thresh, cv2.MORPH_CLOSE, kernel)
            
            return cleaned
            
        except Exception as e:
            # If preprocessing fails, return original image
            return image
    
    async def parse_powerpoint_from_bytes(self, file_bytes: bytes, filename: str) -> Dict[str, Any]:
        """
        Parse a PowerPoint file from bytes
        
        Args:
            file_bytes: PowerPoint file content as bytes
            filename: Original filename
            
        Returns:
            Dictionary containing extracted text and metadata
        """
        # Create temporary file
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
            temp_file.write(file_bytes)
            temp_file_path = temp_file.name
        
        try:
            # Parse the temporary file
            result = await self.parse_powerpoint(temp_file_path)
            result["file_name"] = filename
            return result
        finally:
            # Clean up temporary file
            if os.path.exists(temp_file_path):
                os.unlink(temp_file_path)
    
    def get_supported_formats(self) -> List[str]:
        """Get list of supported file formats"""
        return ['.pptx', '.ppt']
    
    def validate_file(self, file_path: str) -> Tuple[bool, str]:
        """
        Validate if file is a supported PowerPoint format
        
        Args:
            file_path: Path to the file
            
        Returns:
            Tuple of (is_valid, error_message)
        """
        if not os.path.exists(file_path):
            return False, "File does not exist"
        
        file_ext = Path(file_path).suffix.lower()
        if file_ext not in self.get_supported_formats():
            return False, f"Unsupported file format: {file_ext}"
        
        try:
            # Try to open the file to validate it's a valid PowerPoint
            Presentation(file_path)
            return True, ""
        except Exception as e:
            return False, f"Invalid PowerPoint file: {str(e)}"
